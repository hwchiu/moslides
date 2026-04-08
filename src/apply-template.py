#!/usr/bin/env python3
"""
Apply a .pptx master template to a content-only .pptx.

Strategy:
  1. Open the content pptx (10" x 5.5" from pptxgenjs).
  2. Create a new presentation using the template as its slide master.
  3. For each content slide, add a blank slide (template's Blank layout),
     then copy all shapes into it — scaling every position/size from the
     source canvas (10 x 5.5") to the target canvas (13.33 x 7.5").
  4. Save the result.

Usage:
  python3 src/apply-template.py output/part1.pptx templates/template.pptx output/part1_templated.pptx
"""

import sys
import copy
import lxml.etree as etree
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn

# ── EMU constants ────────────────────────────────────────────────────────────
EMU_PER_INCH = 914400

SRC_W = 10.0    # source canvas width  (inches)
SRC_H = 5.5     # source canvas height (inches)


def scale_xfrm(xfrm_el, sx, sy):
    """Scale <a:xfrm> off/ext/chOff/chExt in-place."""
    for tag, xscale, yscale in [
        (qn("a:off"),   sx, sy),
        (qn("a:ext"),   sx, sy),
        (qn("a:chOff"), sx, sy),
        (qn("a:chExt"), sx, sy),
    ]:
        el = xfrm_el.find(tag)
        if el is None:
            continue
        if "x" in el.attrib:
            el.attrib["x"] = str(int(round(int(el.attrib["x"]) * xscale)))
        if "y" in el.attrib:
            el.attrib["y"] = str(int(round(int(el.attrib["y"]) * yscale)))
        if "cx" in el.attrib:
            el.attrib["cx"] = str(int(round(int(el.attrib["cx"]) * xscale)))
        if "cy" in el.attrib:
            el.attrib["cy"] = str(int(round(int(el.attrib["cy"]) * yscale)))


def scale_shape_tree(spTree, sx, sy):
    """Recursively scale every <a:xfrm> found inside spTree."""
    for xfrm in spTree.iter(qn("a:xfrm")):
        scale_xfrm(xfrm, sx, sy)


def scale_font_sizes(spTree, scale):
    """Scale font sizes (<a:sz> values are in hundredths of a point)."""
    for sz_el in spTree.iter(qn("a:sz")):
        try:
            val = int(sz_el.attrib["val"])
            sz_el.attrib["val"] = str(int(round(val * scale)))
        except (KeyError, ValueError):
            pass


def apply_template(content_path, template_path, output_path):
    content  = Presentation(content_path)
    template = Presentation(template_path)

    # Derive target canvas from template
    tgt_w_emu = template.slide_width
    tgt_h_emu = template.slide_height
    tgt_w = tgt_w_emu / EMU_PER_INCH
    tgt_h = tgt_h_emu / EMU_PER_INCH

    src_w_emu = content.slide_width
    src_h_emu = content.slide_height

    sx = tgt_w_emu / src_w_emu   # EMU scale factor X
    sy = tgt_h_emu / src_h_emu   # EMU scale factor Y
    font_scale = (sx + sy) / 2   # proportional font scale (average)

    print(f"Source : {src_w_emu/EMU_PER_INCH:.2f}\" x {src_h_emu/EMU_PER_INCH:.2f}\"")
    print(f"Target : {tgt_w:.2f}\" x {tgt_h:.2f}\"")
    print(f"Scale  : sx={sx:.4f}  sy={sy:.4f}  font={font_scale:.4f}")

    # Remove all pre-existing slides from the template presentation (clean OOXML removal)
    from pptx.oxml.ns import qn as _qn
    prs_part = template.part
    sldIdLst = template.element.find(_qn("p:sldIdLst"))
    for sldId in list(sldIdLst):
        rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        if rId:
            try:
                prs_part.drop_rel(rId)
            except Exception:
                pass
        sldIdLst.remove(sldId)

    # Use "Blank" layout (idx 6) so the master background shows through
    blank_layout = template.slide_layouts[6]  # "Blank"

    for slide_idx, src_slide in enumerate(content.slides):
        new_slide = template.slides.add_slide(blank_layout)
        dst_spTree = new_slide.shapes._spTree

        # Copy every sp/pic/graphicFrame from the source slide
        src_spTree = src_slide.shapes._spTree
        for child in src_spTree:
            tag = child.tag.split("}")[-1]
            if tag in ("nvGrpSpPr", "grpSpPr"):
                continue  # skip group metadata
            node = copy.deepcopy(child)
            scale_shape_tree(node, sx, sy)
            scale_font_sizes(node, font_scale)
            dst_spTree.append(node)

        print(f"  slide {slide_idx + 1:02d} → copied {len(src_spTree) - 2} shapes")

    template.save(output_path)
    print(f"\n✅  Saved: {output_path}")


if __name__ == "__main__":
    if len(sys.argv) != 4:
        print("Usage: python3 apply-template.py <content.pptx> <template.pptx> <output.pptx>")
        sys.exit(1)
    apply_template(sys.argv[1], sys.argv[2], sys.argv[3])
